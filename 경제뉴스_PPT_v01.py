import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from email.utils import parsedate_to_datetime
import urllib.parse
import os
import threading
import webbrowser
import json
import glob
from datetime import date

# ==========================================
# Google News RSS 검색 설정
# ==========================================
GOOGLE_NEWS_RSS = "https://news.google.com/rss/search"
HEADERS = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 '
                  '(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
}


class NewsToPPTApp:
    def __init__(self, root):
        self.root = root
        self.root.title("경제 뉴스 → PPT 변환기 v1.0")
        self.root.geometry("780x720")
        self.root.configure(bg="#F5F6F7")

        self.articles = []
        self.session_start_tokens = None
        self.setup_ui()
        self.update_token_monitor()

    # ------------------------------------------
    # UI 구성
    # ------------------------------------------
    def setup_ui(self):
        header = tk.Frame(self.root, bg="#2C3E50", height=60)
        header.pack(fill="x")
        header.pack_propagate(False)
        tk.Label(header, text="📰 경제 뉴스 → PPT 변환기",
                 bg="#2C3E50", fg="white",
                 font=("맑은 고딕", 14, "bold")).pack(expand=True)

        monitor = tk.Frame(self.root, bg="#1B2631", height=32)
        monitor.pack(fill="x")
        monitor.pack_propagate(False)
        tk.Label(monitor, text="🪙 Claude 토큰 (오늘):",
                 bg="#1B2631", fg="#F39C12",
                 font=("맑은 고딕", 9, "bold")).pack(side="left", padx=10)
        self.token_label = tk.Label(monitor,
                 text="로딩 중...",
                 bg="#1B2631", fg="#ECF0F1",
                 font=("Consolas", 9))
        self.token_label.pack(side="left")
        self.token_session_label = tk.Label(monitor,
                 text="",
                 bg="#1B2631", fg="#85C1E9",
                 font=("Consolas", 9))
        self.token_session_label.pack(side="right", padx=10)

        search_frame = tk.LabelFrame(self.root, text=" 검색 조건 ",
                                     font=("맑은 고딕", 10, "bold"),
                                     bg="white", padx=12, pady=10)
        search_frame.pack(fill="x", padx=15, pady=10)

        kw_frame = tk.Frame(search_frame, bg="white")
        kw_frame.pack(fill="x", pady=4)
        tk.Label(kw_frame, text="키워드:", font=("맑은 고딕", 10),
                 bg="white", width=10, anchor="w").pack(side="left")
        self.keyword_entry = tk.Entry(kw_frame, font=("맑은 고딕", 10), width=45)
        self.keyword_entry.pack(side="left", padx=5)
        self.keyword_entry.insert(0, "환율")
        self.keyword_entry.bind("<Return>", lambda e: self.start_search())

        cnt_frame = tk.Frame(search_frame, bg="white")
        cnt_frame.pack(fill="x", pady=4)
        tk.Label(cnt_frame, text="기사 수:", font=("맑은 고딕", 10),
                 bg="white", width=10, anchor="w").pack(side="left")
        self.count_var = tk.IntVar(value=10)
        ttk.Combobox(cnt_frame, textvariable=self.count_var,
                     values=[5, 10, 15, 20, 30], width=8,
                     state="readonly").pack(side="left", padx=5)

        sort_frame = tk.Frame(search_frame, bg="white")
        sort_frame.pack(fill="x", pady=4)
        tk.Label(sort_frame, text="정렬:", font=("맑은 고딕", 10),
                 bg="white", width=10, anchor="w").pack(side="left")
        self.sort_var = tk.StringVar(value="최신순")
        ttk.Combobox(sort_frame, textvariable=self.sort_var,
                     values=["최신순", "관련도순"], width=10,
                     state="readonly").pack(side="left", padx=5)

        btn_frame = tk.Frame(search_frame, bg="white")
        btn_frame.pack(fill="x", pady=8)
        self.search_btn = tk.Button(btn_frame, text="🔍  뉴스 검색",
                                    command=self.start_search,
                                    bg="#3498DB", fg="white",
                                    font=("맑은 고딕", 10, "bold"),
                                    relief="flat", padx=20, pady=6)
        self.search_btn.pack(side="left", padx=5)

        self.ppt_btn = tk.Button(btn_frame, text="📊  PPT 저장",
                                 command=self.save_to_ppt,
                                 bg="#27AE60", fg="white",
                                 font=("맑은 고딕", 10, "bold"),
                                 relief="flat", padx=20, pady=6,
                                 state="disabled")
        self.ppt_btn.pack(side="left", padx=5)

        self.status_label = tk.Label(self.root,
                                     text="키워드를 입력하고 검색하세요.",
                                     font=("맑은 고딕", 9), bg="#F5F6F7",
                                     fg="#555555")
        self.status_label.pack(pady=3)

        result_frame = tk.LabelFrame(self.root, text=" 검색 결과 ",
                                     font=("맑은 고딕", 10, "bold"),
                                     bg="white", padx=10, pady=5)
        result_frame.pack(fill="both", expand=True, padx=15, pady=10)

        self.result_text = tk.Text(result_frame, font=("맑은 고딕", 9),
                                   bg="#FDFEFE", wrap="word", padx=10, pady=8,
                                   relief="flat")
        scrollbar = tk.Scrollbar(result_frame, command=self.result_text.yview)
        self.result_text.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side="right", fill="y")
        self.result_text.pack(side="left", fill="both", expand=True)

        self.result_text.tag_configure("title",
                                       font=("맑은 고딕", 10, "bold"),
                                       foreground="#2980B9")
        self.result_text.tag_configure("source", foreground="#7F8C8D",
                                       font=("맑은 고딕", 8))
        self.result_text.tag_configure("summary", foreground="#34495E")
        self.result_text.tag_configure("idx", foreground="#E67E22",
                                       font=("맑은 고딕", 10, "bold"))

    # ------------------------------------------
    # 검색 처리
    # ------------------------------------------
    def start_search(self):
        keyword = self.keyword_entry.get().strip()
        if not keyword:
            messagebox.showwarning("입력 오류", "검색 키워드를 입력하세요.")
            return

        self.search_btn.config(state="disabled")
        self.ppt_btn.config(state="disabled")
        self.status_label.config(text=f"'{keyword}' 검색 중...")
        self.result_text.delete("1.0", tk.END)

        threading.Thread(target=self.do_search,
                         args=(keyword,), daemon=True).start()

    def do_search(self, keyword):
        try:
            count = self.count_var.get()
            self.articles = self.search_google_news(keyword, count)
            self.root.after(0, self.display_results)
        except Exception as e:
            err = str(e)
            self.root.after(0, lambda: self.show_error(err))

    def search_google_news(self, keyword, count):
        query = urllib.parse.quote(keyword)
        url = f"{GOOGLE_NEWS_RSS}?q={query}&hl=ko&gl=KR&ceid=KR:ko"
        r = requests.get(url, headers=HEADERS, timeout=10)
        r.raise_for_status()

        soup = BeautifulSoup(r.text, 'xml')
        items = soup.select('item')

        articles = []
        for item in items[:count]:
            title_tag = item.select_one('title')
            link_tag = item.select_one('link')
            pubdate_tag = item.select_one('pubDate')
            source_tag = item.select_one('source')
            desc_tag = item.select_one('description')

            title_raw = title_tag.get_text(strip=True) if title_tag else ''
            title, source_from_title = self._split_title_source(title_raw)
            source = (source_tag.get_text(strip=True)
                      if source_tag else source_from_title)

            pub_str = pubdate_tag.get_text(strip=True) if pubdate_tag else ''
            dt, date_display = self._parse_pubdate(pub_str)

            articles.append({
                'title': title,
                'link': link_tag.get_text(strip=True) if link_tag else '',
                'summary': self._clean_description(
                    desc_tag.get_text() if desc_tag else '', title),
                'source': source or '미상',
                'date': date_display,
                '_dt': dt,
            })

        if self.sort_var.get() == "최신순":
            articles.sort(key=lambda a: a['_dt'].timestamp() if a['_dt'] else 0,
                          reverse=True)
        return articles

    def _split_title_source(self, title_raw):
        if ' - ' in title_raw:
            parts = title_raw.rsplit(' - ', 1)
            return parts[0].strip(), parts[1].strip()
        return title_raw, ''

    def _parse_pubdate(self, pub_str):
        if not pub_str:
            return None, ''
        try:
            dt = parsedate_to_datetime(pub_str)
            return dt, dt.strftime("%Y-%m-%d %H:%M")
        except Exception:
            return None, pub_str

    def _clean_description(self, desc_html, title):
        if not desc_html:
            return ''
        inner = BeautifulSoup(desc_html, 'html.parser')
        texts = [t.strip() for t in inner.stripped_strings
                 if t.strip() and t.strip() != title]
        seen = set()
        unique = []
        for t in texts:
            if t in seen:
                continue
            seen.add(t)
            unique.append(t)
        return ' | '.join(unique[:5])

    def display_results(self):
        if not self.articles:
            self.status_label.config(text="검색 결과가 없습니다.")
            self.result_text.insert(tk.END,
                "\n   검색 결과가 없습니다.\n   다른 키워드로 시도해 보세요.")
            self.search_btn.config(state="normal")
            return

        self.status_label.config(
            text=f"✅ {len(self.articles)}개 기사 검색 완료. PPT 저장 가능합니다.")

        for i, art in enumerate(self.articles, 1):
            link_tag = f"link_{i}"
            self.result_text.tag_configure(link_tag,
                                           foreground="#2980B9",
                                           underline=True,
                                           font=("맑은 고딕", 10, "bold"))
            url = art.get('link', '')
            if url:
                self.result_text.tag_bind(link_tag, "<Button-1>",
                    lambda e, u=url: webbrowser.open(u))
                self.result_text.tag_bind(link_tag, "<Enter>",
                    lambda e: self.result_text.config(cursor="hand2"))
                self.result_text.tag_bind(link_tag, "<Leave>",
                    lambda e: self.result_text.config(cursor=""))

            self.result_text.insert(tk.END, f"\n[{i:02d}] ", "idx")
            self.result_text.insert(tk.END, f"{art['title']}\n", link_tag)
            meta = f"      📰 {art['source']}"
            if art['date']:
                meta += f"   📅 {art['date']}"
            self.result_text.insert(tk.END, meta + "\n", "source")
            if art['summary']:
                self.result_text.insert(tk.END,
                    f"      {art['summary']}\n", "summary")

        self.search_btn.config(state="normal")
        self.ppt_btn.config(state="normal")

    def show_error(self, msg):
        self.status_label.config(text="❌ 오류 발생")
        messagebox.showerror("오류", f"검색 실패:\n{msg}")
        self.search_btn.config(state="normal")

    # ------------------------------------------
    # Claude 토큰 사용량 모니터
    # ------------------------------------------
    def update_token_monitor(self):
        try:
            today_total, total_total = self._aggregate_claude_tokens()

            if self.session_start_tokens is None:
                self.session_start_tokens = total_total
            session_used = total_total - self.session_start_tokens

            self.token_label.config(
                text=f"입력 {self._fmt(today_total['input'])}   "
                     f"출력 {self._fmt(today_total['output'])}   "
                     f"캐시읽기 {self._fmt(today_total['cache_read'])}   "
                     f"캐시생성 {self._fmt(today_total['cache_creation'])}   "
                     f"총 {self._fmt(today_total['total'])}")
            self.token_session_label.config(
                text=f"세션 +{self._fmt(session_used)}")
        except Exception as e:
            self.token_label.config(text=f"(읽기 실패: {str(e)[:40]})")

        self.root.after(5000, self.update_token_monitor)

    def _aggregate_claude_tokens(self):
        today_str = date.today().isoformat()
        today = {'input': 0, 'output': 0,
                 'cache_read': 0, 'cache_creation': 0, 'total': 0}
        all_total = 0

        log_dir = os.path.expanduser("~/.claude/projects")
        pattern = os.path.join(log_dir, "*", "*.jsonl")

        for path in glob.glob(pattern):
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    for line in f:
                        line = line.strip()
                        if not line or '"usage"' not in line:
                            continue
                        try:
                            obj = json.loads(line)
                        except json.JSONDecodeError:
                            continue

                        usage = (obj.get('message', {}) or {}).get('usage')
                        if not isinstance(usage, dict):
                            continue

                        i = int(usage.get('input_tokens', 0) or 0)
                        o = int(usage.get('output_tokens', 0) or 0)
                        cr = int(usage.get('cache_read_input_tokens', 0) or 0)
                        cc = int(usage.get('cache_creation_input_tokens', 0) or 0)
                        sub = i + o + cr + cc
                        all_total += sub

                        ts = obj.get('timestamp', '')
                        if ts.startswith(today_str):
                            today['input'] += i
                            today['output'] += o
                            today['cache_read'] += cr
                            today['cache_creation'] += cc
                            today['total'] += sub
            except (IOError, OSError):
                continue

        return today, all_total

    @staticmethod
    def _fmt(n):
        if n >= 1_000_000:
            return f"{n / 1_000_000:.2f}M"
        if n >= 1_000:
            return f"{n / 1_000:.1f}K"
        return str(n)

    # ------------------------------------------
    # PPT 생성
    # ------------------------------------------
    def save_to_ppt(self):
        if not self.articles:
            return

        keyword = self.keyword_entry.get().strip()
        today = datetime.now().strftime("%Y%m%d")
        default_name = f"경제뉴스_{keyword}_{today}.pptx"

        path = filedialog.asksaveasfilename(
            title="PPT 저장 위치 선택",
            defaultextension=".pptx",
            initialfile=default_name,
            filetypes=[("PowerPoint", "*.pptx")]
        )
        if not path:
            return

        try:
            self.create_ppt(path, keyword)
            if messagebox.askyesno("완료",
                f"PPT 저장 완료:\n{path}\n\n폴더를 열까요?"):
                os.startfile(os.path.dirname(path))
        except Exception as e:
            messagebox.showerror("오류", f"PPT 생성 실패:\n{e}")

    def create_ppt(self, path, keyword):
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.add_summary_slide(slide, keyword)

        prs.save(path)

    def _add_textbox(self, slide, text, left, top, width, height,
                     font_size=18, bold=False, color=(0, 0, 0),
                     align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = RGBColor(*color)
        return box

    def _add_rect(self, slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()
        return shape

    def add_summary_slide(self, slide, keyword):
        self._add_rect(slide, 0, 0, Inches(13.33), Inches(1.1),
                       (0x2C, 0x3E, 0x50))

        self._add_textbox(slide, f"📰  경제 뉴스  —  {keyword}",
                          Inches(0.4), Inches(0.2),
                          Inches(8.5), Inches(0.7),
                          font_size=26, bold=True,
                          color=(0xFF, 0xFF, 0xFF))

        meta_text = (f"{datetime.now().strftime('%Y-%m-%d')}"
                     f"   |   총 {len(self.articles)}건")
        self._add_textbox(slide, meta_text,
                          Inches(9.0), Inches(0.4),
                          Inches(4.0), Inches(0.5),
                          font_size=14, bold=False,
                          color=(0xEC, 0xF0, 0xF1),
                          align=PP_ALIGN.RIGHT)

        self._add_rect(slide, 0, Inches(1.1),
                       Inches(13.33), Inches(0.04),
                       (0xE6, 0x7E, 0x22))

        count = max(len(self.articles), 1)
        available_h = 6.1
        row_h = min(0.9, available_h / count)

        if row_h >= 0.70:
            title_pt, meta_pt = 18, 12
        elif row_h >= 0.45:
            title_pt, meta_pt = 14, 10
        elif row_h >= 0.32:
            title_pt, meta_pt = 12, 9
        elif row_h >= 0.24:
            title_pt, meta_pt = 10, 8
        else:
            title_pt, meta_pt = 9, 7

        top_start = 1.3
        for i, art in enumerate(self.articles):
            top = Inches(top_start + row_h * i)
            self._add_article_row(slide, art, i + 1, top,
                                  Inches(row_h), title_pt, meta_pt)

    def _add_article_row(self, slide, art, idx, top, height,
                         title_pt, meta_pt):
        box = slide.shapes.add_textbox(Inches(0.4), top,
                                       Inches(12.5), height)
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        r1 = p.add_run()
        r1.text = f"[{idx:02d}]  "
        r1.font.size = Pt(title_pt)
        r1.font.bold = True
        r1.font.name = "맑은 고딕"
        r1.font.color.rgb = RGBColor(0xE6, 0x7E, 0x22)

        r2 = p.add_run()
        r2.text = art['title']
        r2.font.size = Pt(title_pt)
        r2.font.bold = True
        r2.font.name = "맑은 고딕"
        r2.font.color.rgb = RGBColor(0x29, 0x80, 0xB9)
        r2.font.underline = True
        if art['link']:
            r2.hyperlink.address = art['link']

        meta = f"   ({art['source']}"
        if art['date']:
            meta += f" · {art['date']}"
        meta += ")"
        r3 = p.add_run()
        r3.text = meta
        r3.font.size = Pt(meta_pt)
        r3.font.name = "맑은 고딕"
        r3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)


if __name__ == "__main__":
    root = tk.Tk()
    app = NewsToPPTApp(root)
    root.mainloop()
